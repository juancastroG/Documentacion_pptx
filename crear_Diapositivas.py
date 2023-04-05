from pptx import Presentation

prs = Presentation()#Creamos el objeto presentacion que contiene todos los metodos de las diapositivas

for i in range(0,10): #Creamos un ciclo para recorrer los 10 primeros estilos de diapositivas que existen
    # 0: Es la hoja principal (Titulo general y un subtitulo)
    # 1: Tituo y texto
    # 6: Hoja en blanco
    diapositiva = prs.slide_layouts[i] #Creamos un estilo de diapositiva
    presentacion =prs.slides.add_slide(diapositiva) #Le agregamos a nuestra presentacion las diferentes diapositivas

prs.save("diapositivas.pptx")#Guardamos la presentacion en un archivo con un nombre.
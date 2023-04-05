from pptx import Presentation
# Crear portada 
#--------------------------------------------
def portada(tt: str, subtt: str):
    
    prs = Presentation() # Creamos el objeto presentacion
    title_slide_layout = prs.slide_layouts[0] # Creamos la diapositiva 0 para la portada
    slide = prs.slides.add_slide(title_slide_layout) # añadimos la diapositiva a la presentacion
    title = slide.shapes.title #Obtenemos el objeto titulo
    subtitle = slide.placeholders[1] #Obtenemos el cuadro de texto del subtitulo

    title.text = tt #Asignamos un titulo
    subtitle.text = subtt #Asignamos un subtitulo

    prs.save('automatizacion.pptx')

#----------------------------------------------
#Creamos una diapostiva con texto y sublistas

def texto_simple(tt, txt1, txt2, txt3):
    prs = Presentation() # Creamos el objeto presentacion
    bullet_slide_layout = prs.slide_layouts[1]  # Creamos la diapositiva 1 para la portada
    slide = prs.slides.add_slide(bullet_slide_layout) # añadimos la diapositiva a la presentacion
    shapes = slide.shapes #Obtenemos todos los objetos interactivos de la diapositiva
    title_shape = shapes.title #Obtenemos el objeto titulo
    body_shape = shapes.placeholders[1] #Obtenemos el espacio del cuadro de texto

    title_shape.text = tt #Asignamos un titulo
    tf = body_shape.text_frame # Obtenemos el objeto cuadro de texto 
    tf.text = txt1 #Agregamos un texto al cuadro de texto

    p = tf.add_paragraph() #añadimos un parrafo al cuadro de texto
    p.text = txt2 # Agregamos texto al parrafo recien creado
    p.level = 1 #Se usa para crear sublistas con formato

    p = tf.add_paragraph() #Añadimos otro parrafo en el cuadro de texto
    p.text = txt3 # Agregamos texto al parrafo recien creado
    p.level = 2 #Creamos una sub-sublista 

    prs.save('test.pptx') #Guardamos el archivo

#-----------------------------------
#Creamos una diapositiva con texto y cambiamos las propiedades

def texto_propiedades(tt, tx1, tx2):
    from pptx.util import Inches, Pt #Importamos las medidas de pulgadas y el de letra

    prs = Presentation() # Creamos el objeto presentacion
    blank_slide_layout = prs.slide_layouts[6]  # Creamos la diapositiva 1 para la portada
    slide = prs.slides.add_slide(blank_slide_layout) # añadimos la diapositiva a la presentacion

    left = top = width = height = Inches(1) #Asignamos las medidas y margenes que se usaran 
    txBox = slide.shapes.add_textbox(left, top, width, height) #Creamos un cuadro de texto y le ponemos margenes y las medidas 
    tf = txBox.text_frame 

    tf.text = tt #Agregamos texto

    p = tf.add_paragraph() #Creamos una parrafo 
    p.text = tx1 #Añadimos texto a este parrafo
    p.font.bold = True #Accedemos a la propiedad de negrilla y la activamos

    p = tf.add_paragraph() #Creamos un parrafo
    p.text = tx2 #Agregamos texto a ese parrafo
    p.font.size = Pt(40) #Cambiamos la propiedad del tamaño de letra

    prs.save('test.pptx') #Guardamos el archivo

#--------------------------------------------------
# Creamos una diapostiva y anexamos 2 imagenes

def imagenes(path):

    from pptx.util import Inches #Importamos las medidas de pulgadas

    img_path = path #Agregamos la ruta de la imagen

    prs = Presentation() # Creamos el objeto presentacion
    blank_slide_layout = prs.slide_layouts[6] # Creamos la diapositiva 1 para la portada
    slide = prs.slides.add_slide(blank_slide_layout) # añadimos la diapositiva a la presentacion

    left = top = Inches(1) #Creamos las margenes para la imagen
    pic = slide.shapes.add_picture(img_path, left, top) #Añadimos la imagen a la diapositiva

    left = Inches(5) #Creamos otras margenes para una segunda imagen
    height = Inches(5.5)
    pic = slide.shapes.add_picture(img_path, left, top, height=height) #Añadimos esa segunda imagen a la diapostiva

    prs.save('test.pptx') #Guardamos elarchivo
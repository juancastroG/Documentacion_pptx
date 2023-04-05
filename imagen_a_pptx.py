from pptx import Presentation
from pptx.util import Inches #Importamos la medicion en pulgadas 

img_path = "pico.jpg" #Damos la ruta de la imagen

prs = Presentation() #Creamos el objeto presentacion
blank_slide_layout = prs.slide_layouts[6] #Creamos una diapostiva en blanco
slide = prs.slides.add_slide(blank_slide_layout) #Agregamos la diapostiva a la presentacion
left = top = Inches(2) #Definimos la margen izq y arriba de la imagen en pulgadas

#pic = slide.shapes.add_picture(img_path, left, top) #Agregamos la imagen con los parametros basicos 

#pic = slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.75), #Agregamos la imagen cambiando las medidas de la imagen
#width=Inches(9), height=Inches(5))

pic = slide.shapes.add_picture(img_path, left, top, Inches(3)) #Agregamos la imagen especificando cuanto queremos que mida, 
#El mismo lo ajustara para que no pierda relacion de altura y anchura
prs.save("imagen.pptx") #Guardamos la presentacion


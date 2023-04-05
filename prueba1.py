from pptx import Presentation #Importamos el modulo pptx
import os #Importamos os para poder ejecutar el archivo una vez creado
 
 
#Creamos una funcion que nos permite crear una presentacion, sencilla con titulos y texto 
def create_presentation(filename, content): 
    
    content2 = list() #Creamos una lista donde almacenaremos el el texto de las diferentes diapositivas
    for line in content.split("\n\n"): # Dividimos el texto de entrada por 'parrafos'
        line = line.split("\n") # creamos una lista separa por un salto de linea, para diferencias titulo y contenido
        content2.append(line) #Lo añadimos a la lista
    print(content2) 
    prs = Presentation() #Creamos el objeto prs que contiene todos los items de una presentación
    layout = prs.slide_layouts[1] #Obtenemos un layout el '[1]' hace referencia a una diapostiva basica de un titulo y un cuadro de contenido
    for lst in content2: #Recorremos la lista
        slide = prs.slides #Creamos una diapositiva
        slide = slide.add_slide(layout) #Le añadimos el layout previamente creado 
        slide.shapes.title.text = lst[0] #Agregamos el titulo almacenado en nuestra lista
        body = slide.placeholders[1] #Creamos la margen donde ira el texto
        tf = body.text_frame #Creamos el cuadro de texto en el placeholder
        if "-" in lst[1]: #Verificamos si queremos una lista de items
            for item in lst[1].split("-")[1:]: #recorremos cada uno los items que queremos en la lista 
                p = tf.add_paragraph() #Creamos un parafo en el cuadro de texto
                p.text = item #Añadimos item por item
                p.level = 1 
        else: #En caso de que sea solo texto lo asignamos al cuadro de texto
            tf.text = lst[1] 
    # Guardamos y abrimos la diapositiva
    prs.save(filename)
    os.startfile(filename)
 
 
content = """Capitulo 1
Inicio


Capitulo 2
my lista: - pan - huevos - sal - cafe"""
 
create_presentation("ejemplo.pptx", content)